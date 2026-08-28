#!/usr/bin/env python3
"""
Generate a PowerPoint (.pptx) Pitch Deck for EduStudent Sight
Uses python-pptx to generate a 16:9 widescreen presentation with Catppuccin Mocha dark theme.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Output file path
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "data")
os.makedirs(OUTPUT_DIR, exist_ok=True)
PPTX_PATH = os.path.join(OUTPUT_DIR, "EduStudent_Sight_Presentation.pptx")

# -------------------------------------------------------------
# PALETTE & DESIGN SYSTEM (Catppuccin Mocha Inspired)
# -------------------------------------------------------------
COLOR_BG = RGBColor(30, 30, 46)          # #1E1E2E (Dark base)
COLOR_SURFACE = RGBColor(37, 37, 56)     # #252538 (Card surface)
COLOR_ELEVATED = RGBColor(49, 50, 68)    # #313244 (Elevated elements)
COLOR_TEXT_WHITE = RGBColor(205, 214, 244) # #CDD6F4 (Primary text)
COLOR_TEXT_MUTED = RGBColor(166, 173, 200) # #A6ADC8 (Secondary text)
COLOR_PRIMARY = RGBColor(137, 180, 250)  # #89B4FA (Blue highlight)
COLOR_ACCENT = RGBColor(180, 190, 254)   # #B4BEFE (Lavender accent)
COLOR_SUCCESS = RGBColor(166, 227, 161)  # #A6E3A1 (Green safe)
COLOR_WARNING = RGBColor(249, 226, 175)  # #F9E2AF (Yellow warning)
COLOR_DANGER = RGBColor(243, 139, 168)   # #F38BA8 (Red critical)
COLOR_BORDER = RGBColor(69, 71, 90)      # #45475A (Card border)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen standard
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide

    def set_slide_background(slide, color=COLOR_BG):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, title, category="EDUSTUDENT SIGHT • AI ACADEMIC INTELLIGENCE"):
        # Category tag
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        # Title
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.name = "Calibri"
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # =========================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_BG)

    # Decorative hero card
    add_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), COLOR_SURFACE, COLOR_PRIMARY)

    # Title content
    tb1 = slide1.shapes.add_textbox(Inches(1.3), Inches(1.3), Inches(10.7), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "INSTITUTIONAL AI PLATFORM • VIGNAN UNIVERSITY"
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p = tf1.add_paragraph()
    p.text = "EduStudent Sight"
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_before = Pt(8)

    p = tf1.add_paragraph()
    p.text = "Autonomous Multi-Signal Academic Early-Warning, Telemetry & Intervention Intelligence Platform"
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(12)

    p = tf1.add_paragraph()
    p.text = "Empowering Institutional Leadership, Course Faculty & Mentors with Explainable AI & Closed-Loop Interventions"
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(8)

    # Pill tags
    p = tf1.add_paragraph()
    p.text = "✓ 4-Tier RBAC Architecture   |   ✓ Multi-Model LLM Hub (Gemini, Gemma, Groq)   |   ✓ 100% Explainable Math Engine"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS
    p.space_before = Pt(36)

    # =========================================================
    # SLIDE 2: THE PROBLEM STATEMENT
    # =========================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Institutional Crisis: Why Academic Early Warnings Fail")

    cards_data_s2 = [
        ("1. Data Isolation & Silos", "Attendance records, Continuous Internal Evaluation (CIE) marks, and LMS access logs exist in disjointed systems with no unified correlation.", COLOR_DANGER),
        ("2. Opaque 'Black-Box' Models", "Generic AI solutions spit out risk percentages without mathematical explanations, destroying faculty trust and preventing root-cause fixes.", COLOR_WARNING),
        ("3. Broken Follow-Up Loops", "Even when students are flagged, counseling sessions are informal. Commitments are lost and completion sign-offs lack supervisory oversight.", COLOR_PRIMARY),
        ("4. Delayed Interventions", "At-risk signals are only caught after semester exams when it is too late to remediate low attendance or conceptual bottlenecks.", COLOR_ACCENT),
    ]

    for i, (ctitle, cdesc, ccolor) in enumerate(cards_data_s2):
        top_pos = Inches(1.5 + (i * 1.35))
        add_card(slide2, Inches(0.8), top_pos, Inches(11.733), Inches(1.15), COLOR_SURFACE, ccolor)
        
        tb = slide2.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.1), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.name = "Calibri"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ccolor

        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.name = "Calibri"
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(3)

    # =========================================================
    # SLIDE 3: THE EDUSTUDENT SIGHT SOLUTION
    # =========================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "The Solution: Unified Telemetry & Closed-Loop Action")

    pillars = [
        ("📊 Multi-Signal Ingestion", "Synchronizes attendance %, Continuous Internal Evaluation (CIE), LMS platform footprints, and extracurricular memberships into a unified 360° student vector.", COLOR_PRIMARY),
        ("📐 Explainable Risk Engine", "Deterministic mathematical formula balancing Attendance (40%), Scaled CGPA (35%), and LMS (25%) with fully configurable institutional thresholds.", COLOR_SUCCESS),
        ("🤝 Closed-Loop Reviews", "Two-party verification where mentors submit session commitments and course instructors or admins conduct structured completion sign-offs.", COLOR_WARNING),
        ("🤖 Autonomous AI Studio", "Multi-LLM engine injected with live cohort telemetry for instant student diagnosis, remedial planning, and statistical trend analysis.", COLOR_ACCENT)
    ]

    for i, (ptitle, pdesc, pcolor) in enumerate(pillars):
        col = i % 2
        row = i // 2
        left_pos = Inches(0.8 + (col * 5.95))
        top_pos = Inches(1.5 + (row * 2.7))

        add_card(slide3, left_pos, top_pos, Inches(5.78), Inches(2.45), COLOR_SURFACE, pcolor)
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.2), Inches(5.28), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ptitle
        p.font.name = "Calibri"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = pcolor

        p2 = tf.add_paragraph()
        p2.text = pdesc
        p2.font.name = "Calibri"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(8)

    # =========================================================
    # SLIDE 4: 4-TIER ROLE-BASED ACCESS CONTROL (RBAC)
    # =========================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "4-Tier RBAC Architecture: Tailored Institutional Personas")

    roles = [
        ("🛡️ Administrator", "System Admin, Deans, HODs", "• Complete cohort-wide governance\n• Faculty signup approvals & decline reasons\n• Dynamic risk thresholds & AI key configuration\n• CSV & Multi-tab Excel bulk import/export", COLOR_DANGER),
        ("🎓 Course Faculty", "Subject Instructors, Lab Leads", "• Teaching roster oversight (DBMS, OS, Math)\n• Subject Continuous Internal Evaluation (CIE)\n• Interventions inception for failing students\n• Review & sign-off on mentor completion requests", COLOR_PRIMARY),
        ("🧭 Student Mentor", "Assigned Academic Mentors", "• Mentee risk radar & attendance habit monitoring\n• 1-on-1 counseling scheduling (Cabin / Meet)\n• Submit session outcome commitments\n• Address reviewer revision feedback loops", COLOR_SUCCESS),
        ("👤 Monitored Student", "Enrolled CSE Undergraduates", "• Personal 360° academic telemetry radar\n• Per-subject marks, attendance & rank view\n• Mentoring action items & historical commitments\n• 24/7 AI Academic Assistant Tutor", COLOR_ACCENT)
    ]

    for i, (rtitle, rsub, rbullets, rcolor) in enumerate(roles):
        left_pos = Inches(0.8 + (i * 2.97))
        add_card(slide4, left_pos, Inches(1.5), Inches(2.82), Inches(5.4), COLOR_SURFACE, rcolor)

        tb = slide4.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.65), Inches(2.52), Inches(5.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = rcolor

        p_sub = tf.add_paragraph()
        p_sub.text = rsub
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_sub.space_before = Pt(2)

        p_body = tf.add_paragraph()
        p_body.text = rbullets
        p_body.font.name = "Calibri"
        p_body.font.size = Pt(11.5)
        p_body.font.color.rgb = COLOR_TEXT_WHITE
        p_body.space_before = Pt(10)

    # =========================================================
    # SLIDE 5: EXPLAINABLE AI RISK MATHEMATICS
    # =========================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Explainable AI Engine: Transparent Mathematical Formulation")

    # Formula Box
    add_card(slide5, Inches(0.8), Inches(1.45), Inches(11.733), Inches(2.1), COLOR_ELEVATED, COLOR_ACCENT)
    tb_f = slide5.shapes.add_textbox(Inches(1.1), Inches(1.55), Inches(11.133), Inches(1.9))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    p = tf_f.paragraphs[0]
    p.text = "MATHEMATICAL ENGAGEMENT & RISK SCORE FORMULATION"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p2 = tf_f.add_paragraph()
    p2.text = "Engagement Index (E) = (Attendance × 0.40) + (CGPA_scaled × 0.35) + (LMS_Score × 0.25)\nRisk Score (R) = max( 0, min( 100, 100.0 - Engagement Index ) )"
    p2.font.name = "Courier New"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_before = Pt(6)

    p3 = tf_f.add_paragraph()
    p3.text = "• CGPA_scaled = CGPA × 10.0  |  • Dynamic Cutoffs: Low (≤30%), Moderate (31-59%), High/Critical (≥60%)"
    p3.font.name = "Calibri"
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_before = Pt(6)

    # Comparison Cards
    add_card(slide5, Inches(0.8), Inches(3.75), Inches(5.75), Inches(3.1), COLOR_SURFACE, COLOR_SUCCESS)
    tb_c1 = slide5.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(5.35), Inches(2.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "🟢 Nominal Student (25CS001 — V. Sri Udbhav)"
    p.font.name = "Calibri"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS

    p_body = tf_c1.add_paragraph()
    p_body.text = "• Attendance: 84% (Weight: 40% → 33.6 pts)\n• CGPA: 8.40 (Scaled: 84.0 | Weight: 35% → 29.4 pts)\n• LMS Score: 88% (Weight: 25% → 22.0 pts)\n• Total Engagement (E): 85.0 pts\n• Calculated Risk (R): 15.0% [LOW RISK / SAFE]"
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = COLOR_TEXT_WHITE
    p_body.space_before = Pt(8)

    add_card(slide5, Inches(6.783), Inches(3.75), Inches(5.75), Inches(3.1), COLOR_SURFACE, COLOR_DANGER)
    tb_c2 = slide5.shapes.add_textbox(Inches(7.0), Inches(3.9), Inches(5.35), Inches(2.8))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "🔴 Critical Risk Student (25CS005 — Arjun Patel)"
    p.font.name = "Calibri"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_DANGER

    p_body = tf_c2.add_paragraph()
    p_body.text = "• Attendance: 61% (Weight: 40% → 24.4 pts)\n• CGPA: 6.80 (Scaled: 68.0 | Weight: 35% → 23.8 pts)\n• LMS Score: 50% (Weight: 25% → 12.5 pts)\n• Total Engagement (E): 60.7 pts\n• Calculated Risk (R): 39.3% [HIGH/CRITICAL WARNING]"
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = COLOR_TEXT_WHITE
    p_body.space_before = Pt(8)

    # =========================================================
    # SLIDE 6: CLOSED-LOOP INTERVENTIONS WORKFLOW
    # =========================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Closed-Loop Interventions & Two-Party Review Queue")

    steps = [
        ("Step 1: Detection", "Autonomous radar detects attendance slump (<75%) or internal marks drop.", COLOR_PRIMARY),
        ("Step 2: Inception", "Faculty or Mentor schedules 1-on-1 session with specific remedial venue & agenda.", COLOR_ACCENT),
        ("Step 3: Counseling", "Session conducted. Student makes concrete commitments and submits missing assignments.", COLOR_WARNING),
        ("Step 4: Completion Req", "Mentor documents session outcome and requests formal sign-off from initiator.", COLOR_SUCCESS),
        ("Step 5: Review & Audit", "Initiator approves completion OR rejects with specific revision instructions.", COLOR_DANGER)
    ]

    for i, (stitle, sdesc, scolor) in enumerate(steps):
        left_pos = Inches(0.8 + (i * 2.38))
        add_card(slide6, left_pos, Inches(1.6), Inches(2.25), Inches(5.2), COLOR_SURFACE, scolor)
        
        tb = slide6.shapes.add_textbox(left_pos + Inches(0.12), Inches(1.8), Inches(2.01), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = scolor

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(10)

    # =========================================================
    # SLIDE 7: MULTI-MODEL LLM HUB
    # =========================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Autonomous AI Studio: Multi-LLM Routing & Resilience")

    models_data = [
        ("⚡ Google Gemini 3.5 Flash", "Default flagship model delivering sub-second reasoning with live cohort telemetry injection.", "gemini-3.5-flash-lite, gemini-3.6-flash", COLOR_PRIMARY),
        ("💎 Google Gemma 4 (31B / MoE)", "Google open-weights instruction model executed via standard Gemini API key without extra billing.", "gemma-4-31b-it, gemma-4-26b-a4b-it", COLOR_ACCENT),
        ("🚀 Groq Cloud LPUs", "Ultra-fast inference (500+ tok/s) with custom Cloudflare bypass headers and failover routing.", "groq/compound, openai/gpt-oss-120b", COLOR_SUCCESS),
        ("🌐 OpenRouter Free Tier", "Auto-free smart router dynamically balancing across healthy community model endpoints.", "liquid/lfm-2.5-2.6b, google/gemma-4-31b-it", COLOR_WARNING),
        ("🔒 Local Ollama & GPU Tunnel", "100% private local GPU inference supported via ngrok or Cloudflare tunnel URL.", "qwen2.5:7b, llama3.2, mistral", COLOR_DANGER)
    ]

    for i, (mtitle, mdesc, mfails, mcolor) in enumerate(models_data):
        top_pos = Inches(1.45 + (i * 1.15))
        add_card(slide7, Inches(0.8), top_pos, Inches(11.733), Inches(1.05), COLOR_SURFACE, mcolor)
        
        tb = slide7.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mtitle
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = mcolor

        p2 = tf.add_paragraph()
        p2.text = f"{mdesc}   [Failover: {mfails}]"
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(2)

    # =========================================================
    # SLIDE 8: PRODUCTION DATA SCALE
    # =========================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Production Data Scale & Master Excel Artifacts")

    metrics = [
        ("100", "Monitored Students", "Spanning all 4 risk tiers across 6 engineering branches", COLOR_PRIMARY),
        ("50", "Faculty & Mentors", "30 course faculty + 20 specialized counselors", COLOR_SUCCESS),
        ("10", "System Administrators", "Dean, HODs, Exam Cell, Student Welfare, IQAC", COLOR_DANGER),
        ("500", "Subject Marks", "5 core subjects/student with granular CIE & grades", COLOR_WARNING),
        ("195", "Activity Records", "10 clubs (Hackathons, Robotics, Coding, IEEE)", COLOR_ACCENT),
        ("35", "Live Interventions", "Completed (15), Enquiries (8), Revisions (5)", COLOR_PRIMARY)
    ]

    for i, (num, label, subtext, mcolor) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left_pos = Inches(0.8 + (col * 3.97))
        top_pos = Inches(1.5 + (row * 2.7))

        add_card(slide8, left_pos, top_pos, Inches(3.78), Inches(2.45), COLOR_SURFACE, mcolor)
        tb = slide8.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.2), Inches(3.38), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = mcolor

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.name = "Calibri"
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = subtext
        p3.font.name = "Calibri"
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(4)

    # =========================================================
    # SLIDE 9: DUAL-CLOUD DEPLOYMENT ARCHITECTURE
    # =========================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Dual-Cloud Production Deployment: Netlify + Render")

    # Architecture Columns
    add_card(slide9, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.3), COLOR_SURFACE, COLOR_PRIMARY)
    tb_net = slide9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.35), Inches(4.9))
    tf_net = tb_net.text_frame
    tf_net.word_wrap = True
    
    p = tf_net.paragraphs[0]
    p.text = "🌐 Frontend on Netlify Edge CDN"
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p_body = tf_net.add_paragraph()
    p_body.text = "• Zero-Build Pure Static Architecture (HTML5 + ES6 Modules)\n• Global CDN Edge Caching for sub-50ms page load\n• Native _redirects reverse-proxy routes /api/* directly to Render (Eliminating CORS bottlenecks)\n• Dual-Mode Config (Auto-detects localhost vs live Render)\n• Fully responsive Catppuccin Mocha / Latte theme tokens"
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(12.5)
    p_body.font.color.rgb = COLOR_TEXT_WHITE
    p_body.space_before = Pt(14)

    add_card(slide9, Inches(6.783), Inches(1.5), Inches(5.75), Inches(5.3), COLOR_SURFACE, COLOR_SUCCESS)
    tb_ren = slide9.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.35), Inches(4.9))
    tf_ren = tb_ren.text_frame
    tf_ren.word_wrap = True
    
    p = tf_ren.paragraphs[0]
    p.text = "⚙️ Backend on Render Web Service"
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS

    p_body = tf_ren.add_paragraph()
    p_body.text = "• Python 3.11 + Flask REST Controller with Gunicorn WSGI\n• Automated blueprint provisioning via render.yaml\n• Embedded SQLite with auto-healing schema initialization\n• Dual-Sync storage: Settings saved in DB and .env in real-time\n• Multi-candidate LLM router with automatic failover loops\n• 120s extended timeout for complex generative AI reasoning"
    p_body.font.name = "Calibri"
    p_body.font.size = Pt(12.5)
    p_body.font.color.rgb = COLOR_TEXT_WHITE
    p_body.space_before = Pt(14)

    # =========================================================
    # SLIDE 10: DEMO CREDENTIALS & EVALUATION MATRIX
    # =========================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Live Demo Credentials & Evaluation Access Matrix")

    # Table Card
    add_card(slide10, Inches(0.8), Inches(1.45), Inches(11.733), Inches(5.4), COLOR_SURFACE, COLOR_PRIMARY)

    rows_data = [
        ("🛡️ Administrator", "admin", "admin123", "Full system governance, AI settings sync, faculty approval queue, CSV/Excel import/export"),
        ("🛡️ Dean Academics", "ADM001", "ADM001", "Dr. P. Venkateswarlu — Executive cohort early-warning intelligence & institutional PDF audit reports"),
        ("🎓 Course Faculty", "FAC001", "FAC001", "Dr. Ramesh Kumar — DBMS & OS course grading, student CIE marks, review mentor completion requests"),
        ("🧭 Senior Mentor", "MEN001", "MEN001", "Prof. Sunitha Devi — Mentee risk radar, 1-on-1 counseling scheduling, submit session commitments"),
        ("👤 Nominal Student", "25CS001", "25CS001", "V. Sri Udbhav — Top performer (CGPA 8.4, Risk 15%), 360° telemetry radar, 24/7 AI tutor assistant"),
        ("⚠️ Critical Student", "25CS005", "25CS005", "Arjun Patel — Critical risk (CGPA 6.8, Attd 61%, Risk 72%), alert indicators, counseling history")
    ]

    tb_t = slide10.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(5.0))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    for i, (r_role, r_id, r_pw, r_desc) in enumerate(rows_data):
        p = tf_t.paragraphs[0] if i == 0 else tf_t.add_paragraph()
        p.text = f"{r_role}   |   ID: {r_id}   |   Password: {r_pw}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT if "Admin" in r_role or "Dean" in r_role else (COLOR_PRIMARY if "Faculty" in r_role else (COLOR_SUCCESS if "Mentor" in r_role else COLOR_WARNING))
        if i > 0:
            p.space_before = Pt(10)

        p2 = tf_t.add_paragraph()
        p2.text = f"Scope: {r_desc}"
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(2)

    # =========================================================
    # SLIDE 11: INSTITUTIONAL IMPACT & ROADMAP
    # =========================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Institutional Value Proposition & Future Roadmap")

    impacts = [
        ("📈 Measurable Retention Uptick", "Reduces semester dropouts by catching attendance and CIE slumps 4 to 6 weeks before mid-term exams.", COLOR_PRIMARY),
        ("📋 NAAC & NBA Accreditation Ready", "Automates Outcome-Based Education (OBE) course attainment calculations and maintains auditable intervention trails.", COLOR_SUCCESS),
        ("📱 Mobile & WhatsApp Bridge", "Roadmap: Automated SMS and WhatsApp progress alerts dispatched to parents upon attendance drops.", COLOR_WARNING),
        ("🧠 Fine-Tuned Institutional LLM", "Roadmap: Training custom on-premise Gemma 4 weights on syllabus questions, previous papers, and academic regulations.", COLOR_ACCENT)
    ]

    for i, (ititle, idesc, icolor) in enumerate(impacts):
        top_pos = Inches(1.5 + (i * 1.35))
        add_card(slide11, Inches(0.8), top_pos, Inches(11.733), Inches(1.15), COLOR_SURFACE, icolor)
        
        tb = slide11.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.1), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ititle
        p.font.name = "Calibri"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = icolor

        p2 = tf.add_paragraph()
        p2.text = idesc
        p2.font.name = "Calibri"
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = COLOR_TEXT_WHITE
        p2.space_before = Pt(3)

    # =========================================================
    # SLIDE 12: CONCLUSION & Q&A
    # =========================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)

    # Decorative hero card
    add_card(slide12, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), COLOR_SURFACE, COLOR_PRIMARY)

    tb12 = slide12.shapes.add_textbox(Inches(1.3), Inches(1.4), Inches(10.7), Inches(4.7))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = "Calibri"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE

    p = tf12.add_paragraph()
    p.text = "EduStudent Sight is Live & Ready for Evaluation"
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_PRIMARY
    p.space_before = Pt(8)

    p = tf12.add_paragraph()
    p.text = "• Live Presentation Demo: http://127.0.0.1:8080\n• Backend API: http://127.0.0.1:5000\n• GitHub Repository: https://github.com/SriUdbhav/EduStudentSight\n• Master Excel Database: data/EduStudent_Sight_Master_Database.xlsx"
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_before = Pt(20)

    p = tf12.add_paragraph()
    p.text = "Vignan University Department of Computer Science & Engineering • Hackathon 2026"
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(30)

    # Save presentation
    prs.save(PPTX_PATH)
    print(f"✓ PowerPoint presentation generated successfully at: {PPTX_PATH}")
    print(f"✓ Slide Count: {len(prs.slides)} slides (16:9 Widescreen)")

if __name__ == "__main__":
    create_presentation()
